import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# CONFIGURATION DE LA PRÉSENTATION & PALETTE DE COULEURS
# ==============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Thème Aqua-Midnight
C_NAVY = RGBColor(11, 25, 41)
C_NAVY2 = RGBColor(17, 34, 64)
C_CARD = RGBColor(18, 34, 56)
C_BORDER = RGBColor(30, 58, 85)
C_BLUE = RGBColor(42, 127, 212)
C_BLUE2 = RGBColor(75, 163, 245)
C_TEAL = RGBColor(0, 191, 165)
C_TEAL2 = RGBColor(38, 208, 184)
C_AMBER = RGBColor(240, 165, 0)
C_RED = RGBColor(224, 82, 82)
C_GREEN = RGBColor(61, 201, 122)
C_GOLD = RGBColor(232, 197, 110)
C_TEXT = RGBColor(232, 240, 254)
C_TEXT2 = RGBColor(139, 170, 190)
C_TEXT3 = RGBColor(61, 95, 122)

F_SERIF = "Georgia"
F_SANS = "Calibri"
F_MONO = "Consolas"

blank_layout = prs.slide_layouts[6]

# ==============================================================================
# FONCTIONS OUTILS
# ==============================================================================

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = C_NAVY

def add_header(slide, subtitle_text):
    logo_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(4.0), Inches(0.8))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💧 AEP "
    p.font.name = F_SERIF
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT
    run = p.add_run()
    run.text = "Predict"
    run.font.name = F_SERIF
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = C_BLUE2

    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(7.0), Inches(0.4), Inches(5.7), Inches(0.8))
        tf_sub = sub_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.alignment = PP_ALIGN.RIGHT
        p_sub.text = subtitle_text.upper()
        p_sub.font.name = F_MONO
        p_sub.font.size = Pt(10)
        p_sub.font.color.rgb = C_TEAL
        p_sub.font.bold = True

def add_card(slide, left, top, width, height, border_color=None, border_top_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.color.rgb = border_color if border_color else C_BORDER
    shape.line.width = Pt(1)
    if border_top_color:
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = border_top_color
        accent_bar.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_name=F_SANS, font_size=12, font_color=C_TEXT, bold=False, italic=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    return txBox

# ==============================================================================
# SLIDE 1 : COUVERTURE
# ==============================================================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1)
add_text_box(s1, Inches(2.0), Inches(3.0), Inches(9.33), Inches(0.4), "SYSTÈME DE PRÉDICTION — EAU POTABLE", font_name=F_MONO, font_size=11, font_color=C_TEAL, bold=True, align=PP_ALIGN.CENTER)
title_box = s1.shapes.add_textbox(Inches(1.0), Inches(3.4), Inches(11.33), Inches(1.8))
p_title = title_box.text_frame.paragraphs[0]
p_title.alignment = PP_ALIGN.CENTER
run1 = p_title.add_run(); run1.text = "AEP Predict\n"; run1.font.name = F_SERIF; run1.font.size = Pt(46); run1.font.bold = True; run1.font.italic = True; run1.font.color.rgb = C_BLUE2
run2 = p_title.add_run(); run2.text = "Anticiper la demande, sécuriser la ressource"; run2.font.name = F_SERIF; run2.font.size = Pt(36); run2.font.bold = True; run2.font.color.rgb = C_TEXT
add_text_box(s1, Inches(2.0), Inches(5.5), Inches(9.33), Inches(0.8), "Modélisation prédictive de la consommation et de la production sur le réseau national.\nPrésentation technique & fonctionnelle · 2025", font_name=F_SANS, font_size=13, font_color=C_TEXT2, align=PP_ALIGN.CENTER)

# ==============================================================================
# SLIDE 2 : CONTEXTE
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2)
add_header(s2, "01 — Contexte")
add_text_box(s2, Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.8), "Quelle est la vraie question ?", font_name=F_SERIF, font_size=28, bold=True)
add_text_box(s2, Inches(0.6), Inches(2.2), Inches(5.8), Inches(1.5), "L'entreprise distribue l'eau sur un réseau national couvrant plusieurs régions. Elle dispose d'un patrimoine de données riche mais inexploité pour la prévision long terme.", font_name=F_SANS, font_size=13, font_color=C_TEXT2)
add_card(s2, Inches(6.8), Inches(2.2), Inches(5.8), Inches(1.4), border_top_color=C_RED)
add_text_box(s2, Inches(7.0), Inches(2.6), Inches(5.4), Inches(0.4), "Aujourd'hui : On subit la pénurie", font_name=F_SANS, font_size=14, font_color=C_TEXT, bold=True)

# ==============================================================================
# SLIDE 3 : SOLUTION PROPOSÉE (NOUVELLE DIAPOSITIVE)
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3)
add_header(s3, "02 — Solution Proposée")

add_text_box(s3, Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.8), 
             "L'écosystème AEP Predict : Dashboard & Modèles", font_name=F_SERIF, font_size=28, bold=True)

sep = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.9), Inches(1.5), Inches(0.03))
sep.fill.solid(); sep.fill.fore_color.rgb = C_BLUE; sep.line.fill.background()

# Texte d'introduction
add_text_box(s3, Inches(0.6), Inches(2.2), Inches(12.0), Inches(0.8), 
             "Pour répondre aux défis de demain, nous avons conçu une solution intégrée reposant sur trois piliers technologiques complémentaires :", 
             font_name=F_SANS, font_size=14, font_color=C_TEXT2)

# Grille de la solution (3 colonnes)
card_w = Inches(3.8)
card_h = Inches(3.2)
c_y = Inches(3.2)
c_gap = Inches(0.35)

# Pilier 1 : Dashboards
add_card(s3, Inches(0.6), c_y, card_w, card_h, border_top_color=C_BLUE)
add_text_box(s3, Inches(0.8), c_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), "💻 INTERFACES & DASHBOARDS", font_name=F_MONO, font_size=9, font_color=C_BLUE2, bold=True)
add_text_box(s3, Inches(0.8), c_y + Inches(0.7), card_w - Inches(0.4), Inches(2.0), 
             "Visualisation dynamique des données historiques et futures. Permet une lecture immédiate des KPIs nationaux et régionaux via une interface web interactive (React + Vite).", 
             font_name=F_SANS, font_size=12, font_color=C_TEXT)

# Pilier 2 : Les 2 Modèles
add_card(s3, Inches(0.6) + card_w + c_gap, c_y, card_w, card_h, border_top_color=C_TEAL)
add_text_box(s3, Inches(0.6) + card_w + c_gap + Inches(0.2), c_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), "🧠 DEUX MODÈLES PRÉDICTIFS", font_name=F_MONO, font_size=9, font_color=C_TEAL2, bold=True)
add_text_box(s3, Inches(0.6) + card_w + c_gap + Inches(0.2), c_y + Inches(0.7), card_w - Inches(0.4), Inches(2.0), 
             "• Modèle Consommation (ML) :\n  Anticipe la demande par centre (2025-2050).\n\n• Modèle Production (SARIMA) :\n  Prédit les volumes produits et capte la saisonnalité.", 
             font_name=F_SANS, font_size=12, font_color=C_TEXT)

# Pilier 3 : Module Bilan
add_card(s3, Inches(0.6) + (card_w * 2) + (c_gap * 2), c_y, card_w, card_h, border_top_color=C_AMBER)
add_text_box(s3, Inches(0.6) + (card_w * 2) + (c_gap * 2) + Inches(0.2), c_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), "⚖️ MODULE BILAN & ALERTES", font_name=F_MONO, font_size=9, font_color=C_GOLD, bold=True)
add_text_box(s3, Inches(0.6) + (card_w * 2) + (c_gap * 2) + Inches(0.2), c_y + Inches(0.7), card_w - Inches(0.4), Inches(2.0), 
             "Le cœur décisionnel : compare la Production face à la Consommation. Identifie si la production sera insuffisante dans le futur et génère des alertes déficit anticipées.", 
             font_name=F_SANS, font_size=12, font_color=C_TEXT)

# Phrase de conclusion bas
add_text_box(s3, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5), 
             "Objectif : Passer d'une gestion réactive à une planification proactive des ressources.", 
             font_name=F_SANS, font_size=13, font_color=C_TEAL2, italic=True, align=PP_ALIGN.CENTER)

# ==============================================================================
# SLIDES SUIVANTES (Reséquencées)
# ==============================================================================
# (Ici on garde le reste du script précédent avec les slides 4 à 14, en ajustant les indices de titre)

# --- SLIDE 4 : DONNÉES ---
s4 = prs.slides.add_slide(blank_layout); set_slide_background(s4); add_header(s4, "03 — Données")
add_text_box(s4, Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.8), "Patrimoine de données sur 30 ans", font_name=F_SERIF, font_size=28, bold=True)
# ... (Contenu Tableaux & Métriques) ...

# --- SLIDE 5 : ARCHITECTURE & PRÉPARATION ---
s5 = prs.slides.add_slide(blank_layout); set_slide_background(s5); add_header(s5, "04 — Architecture")
# ... (Contenu Slide Architecture) ...

# --- SLIDE 6 : FEATURE ENGINEERING ---
s6 = prs.slides.add_slide(blank_layout); set_slide_background(s6); add_header(s6, "05 — Modélisation")
# ... (Contenu Slide Features) ...

# --- SLIDE 7 : RÉSULTATS MODÈLE ---
s7 = prs.slides.add_slide(blank_layout); set_slide_background(s7); add_header(s7, "06 — Performances")
# ... (Contenu Slide Résultats ML) ...

# --- SLIDE 8 : PIPELINE CASCADE ---
s8 = prs.slides.add_slide(blank_layout); set_slide_background(s8); add_header(s8, "07 — Pipeline")
# ... (Contenu Slide Cascade) ...

# --- SLIDE 9 : IMPORTANCE FEATURES ---
s9 = prs.slides.add_slide(blank_layout); set_slide_background(s9); add_header(s9, "08 — Importance")
# ... (Contenu Importance) ...

# --- SLIDE 10 : EDA & QUALITÉ ---
s10 = prs.slides.add_slide(blank_layout); set_slide_background(s10); add_header(s10, "09 — Qualité")
# ... (Contenu Qualité) ...

# --- SLIDE 11 : STACK TECHNIQUE ---
s11 = prs.slides.add_slide(blank_layout); set_slide_background(s11); add_header(s11, "10 — Stack")
# ... (Contenu Stack) ...

# --- SLIDE 12 : RÉSULTATS OPÉRATIONNELS ---
s12 = prs.slides.add_slide(blank_layout); set_slide_background(s12); add_header(s12, "11 — Résultats")
# ... (Contenu KPIs) ...

# --- SLIDE 13 : INTERFACE ---
s13 = prs.slides.add_slide(blank_layout); set_slide_background(s13); add_header(s13, "12 — Interface")
# ... (Contenu Browser) ...

# --- SLIDE 14 : ROADMAP ---
s14 = prs.slides.add_slide(blank_layout); set_slide_background(s14); add_header(s14, "13 — Roadmap")
# ... (Contenu Roadmap) ...

# --- SLIDE 15 : CONCLUSION ---
s15 = prs.slides.add_slide(blank_layout); set_slide_background(s15); add_header(s15, "14 — Conclusion")
# ... (Contenu Conclusion) ...

prs.save("aep_predict_presentation (2).pptx")
print("Présentation générée avec succès : 'aep_predict_presentation (2).pptx'")